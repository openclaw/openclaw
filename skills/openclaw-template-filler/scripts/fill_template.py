#!/usr/bin/env python3
"""
fill_template.py — Template filler for .docx and .pptx files.

Modes:
  --list TEMPLATE              List all placeholders found in the template
  --read-content FILE          Extract text from a .docx or .txt content file
  --fill TEMPLATE              Fill the template
    --output OUTPUT_PATH       Where to save the filled document
    --replacements JSON        JSON string: {"placeholder": "replacement", ...}
    --replacements-file FILE   JSON file with replacements
"""

import sys, json, os, argparse

PLACEHOLDER_MARKERS = [
    "[Buraya", "[buraya", "(Buraya", "(buraya",
    "[Insert", "(Insert", "[TODO", "[PLACEHOLDER",
    "[BURAYA", "[Company", "[Date", "[Author",
]

def looks_like_placeholder(text):
    t = text.strip()
    if not t:
        return False
    for marker in PLACEHOLDER_MARKERS:
        if marker in t:
            return True
    # bracketed or parenthesized short text
    if (t.startswith('[') and t.endswith(']')) or (t.startswith('(') and t.endswith(')')):
        return True
    return False

def list_placeholders(template_path):
    ext = os.path.splitext(template_path)[1].lower()
    if ext == '.docx':
        from docx import Document
        doc = Document(template_path)
        current_heading = "(no heading)"
        for para in doc.paragraphs:
            t = para.text.strip()
            if not t:
                continue
            if para.style.name.startswith('Heading'):
                current_heading = t
            elif looks_like_placeholder(t):
                print(f"[{current_heading}] {repr(t)}")
    elif ext == '.pptx':
        from pptx import Presentation
        prs = Presentation(template_path)
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t and looks_like_placeholder(t):
                        print(f"[Slide {i}] {repr(t)}")
    else:
        with open(template_path, encoding='utf-8') as f:
            for i, line in enumerate(f, 1):
                if looks_like_placeholder(line.strip()):
                    print(f"[Line {i}] {repr(line.strip())}")

def read_content(content_path):
    ext = os.path.splitext(content_path)[1].lower()
    if ext == '.docx':
        from docx import Document
        doc = Document(content_path)
        text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
        print(text)
    else:
        with open(content_path, encoding='utf-8') as f:
            print(f.read())

def replace_in_paragraph(paragraph, replacements):
    full_text = ''.join(run.text for run in paragraph.runs)
    changed = False
    for old, new in replacements.items():
        if old in full_text:
            full_text = full_text.replace(old, new)
            changed = True
    if changed and paragraph.runs:
        paragraph.runs[0].text = full_text
        for run in paragraph.runs[1:]:
            run.text = ''

def fill_docx(template_path, replacements, output_path):
    from docx import Document
    doc = Document(template_path)
    for para in doc.paragraphs:
        replace_in_paragraph(para, replacements)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    replace_in_paragraph(para, replacements)
    doc.save(output_path)
    print(f"Saved: {output_path}")

def fill_pptx(template_path, replacements, output_path):
    from pptx import Presentation
    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full_text = ''.join(run.text for run in para.runs)
                changed = False
                for old, new in replacements.items():
                    if old in full_text:
                        full_text = full_text.replace(old, new)
                        changed = True
                if changed and para.runs:
                    para.runs[0].text = full_text
                    for run in para.runs[1:]:
                        run.text = ''
    prs.save(output_path)
    print(f"Saved: {output_path}")

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--list', metavar='TEMPLATE', help='List placeholders in template')
    parser.add_argument('--read-content', metavar='FILE', help='Extract text from content file')
    parser.add_argument('--fill', metavar='TEMPLATE', help='Fill the template')
    parser.add_argument('--output', metavar='PATH', help='Output path for filled document')
    parser.add_argument('--replacements', metavar='JSON', help='JSON replacements string')
    parser.add_argument('--replacements-file', metavar='FILE', help='JSON replacements file')
    args = parser.parse_args()

    if args.list:
        list_placeholders(os.path.expanduser(args.list))
    elif args.read_content:
        read_content(os.path.expanduser(args.read_content))
    elif args.fill:
        if args.replacements:
            replacements = json.loads(args.replacements)
        elif args.replacements_file:
            with open(os.path.expanduser(args.replacements_file), encoding='utf-8') as f:
                replacements = json.load(f)
        else:
            print("Error: --replacements or --replacements-file required")
            sys.exit(1)

        template_path = os.path.expanduser(args.fill)
        output_path = os.path.expanduser(args.output) if args.output else \
            os.path.splitext(template_path)[0] + '_filled' + os.path.splitext(template_path)[1]

        ext = os.path.splitext(template_path)[1].lower()
        if ext == '.docx':
            fill_docx(template_path, replacements, output_path)
        elif ext == '.pptx':
            fill_pptx(template_path, replacements, output_path)
        else:
            print(f"Unsupported format: {ext}")
            sys.exit(1)
    else:
        parser.print_help()

if __name__ == '__main__':
    main()
